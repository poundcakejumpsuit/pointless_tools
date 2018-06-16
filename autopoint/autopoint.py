import sys, json, os, shutil, smtplib
from PyQt5.QtWidgets import QMainWindow, QApplication, QWidget, QLabel, QPushButton, QAction, QLineEdit, QMessageBox
from PyQt5.QtGui import QIcon
from PyQt5.QtCore import pyqtSlot
from bs4 import BeautifulSoup as bsp
import requests as rq
import pprint, wget
from pptx import Presentation
from pptx.util import Inches
from email.mime.application import MIMEApplication
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.utils import COMMASPACE, formatdate
from os.path import basename

#note for gmail you must generate an application pw
HOST_ADDR = "user@gmail.com"
HOST_PW = "dinosaurpotatoes"


def send_mail(send_from, send_to, subject, text, files=None,
              server="127.0.0.1"):
    assert isinstance(send_to, list)

    msg = MIMEMultipart()
    msg['From'] = send_from
    msg['To'] = COMMASPACE.join(send_to)
    msg['Date'] = formatdate(localtime=True)
    msg['Subject'] = subject

    msg.attach(MIMEText(text))

    for f in files or []:
        with open(f, "rb") as fil:
            part = MIMEApplication(
                fil.read(),
                Name=basename(f)
            )
        # After the file is closed
        part['Content-Disposition'] = 'attachment; filename="{}"'.format(basename(f))
        msg.attach(part)


    smtp = smtplib.SMTP(server, port=587)
    if server != "127.0.0.1":
        smtp.ehlo()
        smtp.starttls()
        smtp.ehlo()
        smtp.login(HOST_ADDR, HOST_PW)
        smtp.send_message(msg)
        del msg
        return
    smtp.sendmail(send_from, send_to, msg.as_string())
    smtp.close()

def fetch_images(query):
    # addr = "https://www.google.com/search?tbm=isch&q=" + query
    addr = "https://www.google.co.in/search?q={}&source=lnms&tbm=isch".format(query)
    s = rq.session()
    r = s.get(addr)
    soup = bsp(r.text, 'html.parser')
    # print(soup)
    imgs = []
    for im in soup.find_all('a', href=True):
        if im.find('img'):
            imgs.append(im)
    # pp = pprint.PrettyPrinter(indent=4)
    # pp.pprint(imgs)
    try:
        os.mkdir('images')
    except OSError:
        shutil.rmtree('images')
        os.mkdir('images')
    count = 0
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    left = top = Inches(1)
    for im in imgs:
        for child in im.children:
            filename = wget.download(child['src'], out='images/image_{}'.format(count))
            count += 1
    for i in range(count):
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture('images/image_{}'.format(i), left, top)

    title = '{}_image.pptx'.format(query)
    prs.save(title)
    return title


class App(QMainWindow):
 
    def __init__(self):
        super().__init__()
        self.title = 'AutoPoint'
        self.left = 10
        self.top = 10
        self.width = 400
        self.height = 240
        self.initUI()
 
    def initUI(self):
        self.setWindowTitle(self.title)
        self.setGeometry(self.left, self.top, self.width, self.height)
 
        # Create textbox for query
        self.textbox = QLineEdit(self)
        self.textbox.setPlaceholderText('Query')
        self.textbox.move(20, 20)
        self.textbox.resize(280,40)

        # Create emailbox
        self.emailbox = QLineEdit(self)
        self.emailbox.setPlaceholderText('Destination email')
        self.emailbox.move(20, 60)
        self.emailbox.resize(280, 40)
 
        # Create a button in the window
        self.button = QPushButton('Generate slides', self)
        self.button.resize(self.button.sizeHint())
        self.button.move(20,180)

        # Create a label to provide feedback
        self.lbl = QLabel("Idle", self)
        self.lbl.resize(self.lbl.sizeHint())
        self.lbl.move(220, 180)
 
        # connect button to function on_click
        self.button.clicked.connect(self.on_click)
        self.show()
 
    @pyqtSlot()
    def on_click(self):
        textboxValue = self.textbox.text()
        self.lbl.setText("searching: {}".format(textboxValue))
        self.lbl.resize(self.lbl.sizeHint())
        ppt_title = fetch_images(textboxValue)
        send_mail(HOST_ADDR, [self.emailbox.text()], "Here's a useless powerpoint", 
                    ppt_title, [ppt_title], server="smtp.gmail.com")
        self.textbox.setText("")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    ex = App()
    sys.exit(app.exec_())